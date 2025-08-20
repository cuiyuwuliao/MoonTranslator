import json
import os
from dataclasses import dataclass
from typing import List, Dict, Any, Tuple
import sys
import re
# PDF: PyMuPDF
import fitz  # PyMuPDF

# XLSX: openpyxl
from openpyxl import load_workbook
from openpyxl.workbook.workbook import Workbook
from openpyxl.worksheet.worksheet import Worksheet

# PPTX: python-pptx
from pptx import Presentation
from pptx.shapes.group import GroupShape
from pptx.shapes.placeholder import PlaceholderPicture
from pptx.shapes.autoshape import Shape
from pptx.shapes.picture import Picture
from pptx.slide import Slide

currentDir = ""
if getattr(sys, 'frozen', False):
    currentDir = os.path.dirname(sys.executable)
else:
    currentDir = os.path.dirname(os.path.abspath(__file__))
# -----------------------
# Utilities
# -----------------------

def _write_json(path: str, data: Any) -> None:
    with open(path, "w", encoding="utf-8") as f:
        json.dump(data, f, ensure_ascii=False, indent=2)

def _contains_chinese(text: str) -> bool:
    # Regular expression pattern for Chinese characters
    pattern = re.compile(r'[\u4e00-\u9fff]')
    return bool(pattern.search(text))

def _read_json(path: str) -> Any:
    with open(path, "r", encoding="utf-8") as f:
        return json.load(f)


def _texts_to_map(texts_list: List[Dict[str, Any]]) -> Dict[str, str]:
    return {item["textId"]: item.get("text", "") for item in texts_list}

def _get_shape_by_path(slide: Slide, path: List[int]):
    """
    Navigate shapes via path of indices (including through group shapes).
    """
    shape = slide.shapes[path[0]]
    for idx in path[1:]:
        if hasattr(shape, "shapes") and isinstance(shape, GroupShape):
            shape = shape.shapes[idx]
        else:
            return None
    return shape
def add_suffix_to_filename(file_path, suffix):
    # Split the file path into directory, base name, and extension
    directory, filename = os.path.split(file_path)
    name, ext = os.path.splitext(filename)
    
    # Create the new filename with the suffix
    new_filename = f"{name}{suffix}{ext}"
    
    # Join the new filename with the directory
    new_file_path = os.path.join(directory, new_filename)
    
    return new_file_path

# -----------------------
# PDF (fitz / PyMuPDF)
# -----------------------
def extractTextsFromPdf(input_path: str) -> None:
    inputFileName, inputFileExtension = os.path.splitext(os.path.basename(input_path))
    extractionFolderPath = os.path.join(os.path.dirname(input_path), f"{inputFileName}_extraction")
    os.makedirs(extractionFolderPath, exist_ok=True)
    out_meta_json = os.path.join(extractionFolderPath, f"{inputFileName}_texts_meta.json")
    out_texts_json = os.path.join(extractionFolderPath, f"{inputFileName}_texts_content.json")
    doc = fitz.open(input_path)
    texts = []
    text_info = []
    tid_counter = 0
    for page_index in range(len(doc)):
        page = doc[page_index]
        pdict = page.get_text("dict")

        for b_idx, block in enumerate(pdict.get("blocks", [])):
            if block.get("type", 0) != 0:
                continue
            for l_idx, line in enumerate(block.get("lines", [])):
                for s_idx, span in enumerate(line.get("spans", [])):
                    text = span.get("text", "")
                    if text is None:
                        text = ""
                    tid = f"pdf_p{page_index}_b{b_idx}_l{l_idx}_s{s_idx}_{tid_counter}"
                    tid_counter += 1
                    texts.append({"textId": tid, "text": text})
                    bbox = span.get("bbox", None)
                    if bbox is None:
                        bbox = line.get("bbox", [0, 0, 0, 0])
                    meta = {
                        "textId": tid,
                        "page": page_index,
                        "blockIndex": b_idx,
                        "lineIndex": l_idx,
                        "spanIndex": s_idx,
                        "bbox": list(map(float, bbox)),
                    }
                    if "font" in span:
                        meta["font"] = span["font"]
                    if "size" in span:
                        meta["size"] = float(span["size"])
                    if "color" in span:
                        meta["color"] = int(span["color"])
                    text_info.append(meta)
    doc.close()
    _write_json(out_texts_json, texts)
    _write_json(out_meta_json, {
        "originalFilePath": os.path.abspath(input_path),
        "textInfo": text_info
    })
    return out_meta_json, out_texts_json


def importTextsToPdf(meta_json: str) -> None:
    texts_json = meta_json.replace("_meta.json", "_content.json")
    meta = _read_json(meta_json)
    texts_map = _texts_to_map(_read_json(texts_json))
    input_path = meta["originalFilePath"]

    doc = fitz.open(input_path)

    per_page_items: Dict[int, List[Dict[str, Any]]] = {}
    for info in meta.get("textInfo", []):
        tid = info["textId"]
        if tid not in texts_map:
            continue
        page_idx = int(info["page"])
        per_page_items.setdefault(page_idx, []).append(info)

    white = (1, 1, 1)
    black = (0, 0, 0)

    for page_idx, items in per_page_items.items():
        page = doc[page_idx]
        for info in items:
            tid = info["textId"]
            new_text = texts_map.get(tid, "")
            x0, y0, x1, y1 = info.get("bbox", [0, 0, 0, 0])
            rect = fitz.Rect(x0, y0, x1, y1)
            
            # Draw white rectangle to cover original text
            page.draw_rect(rect, color=white, fill=white, stroke_opacity=0, fill_opacity=1.0)

            fontname = "helv"
            if _contains_chinese(new_text):
                fontname = "china-ss"
            size = float(info.get("size", 10))
            
            # Calculate text position - using y1 (bottom of original bbox) minus some adjustment
            # The adjustment accounts for font descent to better align with original text
            text_position = fitz.Point(x0, y1 - size * 0.2)  # Adjust the 0.2 factor as needed
            page.insert_text(
                text_position,
                new_text,
                fontname= fontname, 
                fontsize=size,
                color=black,
            )
    outputFileName = add_suffix_to_filename(input_path, "_Ptext")
    try:
        doc.save(outputFileName, incremental=False)
    except ValueError:
        doc.save(outputFileName, incremental=True)
    doc.close()
    return outputFileName

# -----------------------
# XLSX (openpyxl)
# -----------------------

def extractTextsFromXlsx(input_path: str) -> None:
    inputFileName, inputFileExtension = os.path.splitext(os.path.basename(input_path))
    extractionFolderPath = os.path.join(os.path.dirname(input_path), f"{inputFileName}_extraction")
    os.makedirs(extractionFolderPath, exist_ok=True)
    out_meta_json = os.path.join(extractionFolderPath, f"{inputFileName}_texts_meta.json")
    out_texts_json = os.path.join(extractionFolderPath, f"{inputFileName}_texts_content.json")
    
    wb: Workbook = load_workbook(input_path)
    texts = []
    text_info = []
    tid_counter = 0

    for ws in wb.worksheets:
        for row in ws.iter_rows(values_only=False):
            for cell in row:
                val = cell.value
                if val is None:
                    continue
                text = str(val)
                tid = f"xlsx_{ws.title}_{cell.coordinate}_{tid_counter}"
                tid_counter += 1
                texts.append({"textId": tid, "text": text})
                text_info.append({
                    "textId": tid,
                    "sheetName": ws.title,
                    "cell": cell.coordinate
                })

    wb.close()

    _write_json(out_texts_json, texts)
    _write_json(out_meta_json, {
        "originalFilePath": os.path.abspath(input_path),
        "textInfo": text_info
    })
    return out_meta_json, out_texts_json

def importTextsToXlsx(meta_json: str) -> None:
    texts_json = meta_json.replace("_meta.json", "_content.json")
    meta = _read_json(meta_json)
    texts_map = _texts_to_map(_read_json(texts_json))
    input_path = meta["originalFilePath"]

    wb: Workbook = load_workbook(input_path)
    for info in meta.get("textInfo", []):
        tid = info["textId"]
        if tid not in texts_map:
            continue
        sheet = info["sheetName"]
        cell_ref = info["cell"]
        if sheet not in wb.sheetnames:
            continue
        ws: Worksheet = wb[sheet]
        ws[cell_ref].value = texts_map[tid]
    outputFileName = add_suffix_to_filename(input_path, "_Pimage")
    wb.save(outputFileName)
    wb.close()
    return outputFileName

# -----------------------
# PPTX (python-pptx)
# -----------------------

def _iter_shapes_recursive(shape, path_prefix: List[int]) -> List[Tuple[List[int], Any]]:
    out = []
    if hasattr(shape, "shapes") and isinstance(shape, GroupShape):
        for i, subshape in enumerate(shape.shapes):
            out.extend(_iter_shapes_recursive(subshape, path_prefix + [i]))
    else:
        out.append((path_prefix, shape))
    return out

def extractTextsFromPptx(input_path: str) -> None:
    inputFileName, inputFileExtension = os.path.splitext(os.path.basename(input_path))
    extractionFolderPath = os.path.join(os.path.dirname(input_path), f"{inputFileName}_extraction")
    os.makedirs(extractionFolderPath, exist_ok=True)
    out_meta_json = os.path.join(extractionFolderPath, f"{inputFileName}_texts_meta.json")
    out_texts_json = os.path.join(extractionFolderPath, f"{inputFileName}_texts_content.json")
    prs = Presentation(input_path)
    texts = []
    text_info = []
    tid_counter = 0
    for s_idx, slide in enumerate(prs.slides):
        for sh_idx, shape in enumerate(slide.shapes):
            for path, target in _iter_shapes_recursive(shape, [sh_idx]):
                if hasattr(target, "has_text_frame") and target.has_text_frame:
                    tf = target.text_frame
                    for p_idx, paragraph in enumerate(tf.paragraphs):
                        for r_idx, run in enumerate(paragraph.runs):
                            text = run.text or ""
                            tid = f"pptx_s{s_idx}_p{p_idx}_r{r_idx}_{tid_counter}"
                            tid_counter += 1
                            texts.append({"textId": tid, "text": text})
                            text_info.append({
                                "textId": tid,
                                "slideIndex": s_idx,
                                "shapePath": path,
                                "paragraphIndex": p_idx,
                                "runIndex": r_idx
                            })
                # Check for tables
                if target.has_table:
                    table = target.table
                    for row_idx, row in enumerate(table.rows):
                        for cell_idx, cell in enumerate(row.cells):
                            for p_idx, paragraph in enumerate(cell.text_frame.paragraphs):
                                for r_idx, run in enumerate(paragraph.runs):
                                    text = run.text or ""
                                    tid = f"pptx_s{s_idx}_t{row_idx}_c{cell_idx}_p{p_idx}_r{r_idx}_{tid_counter}"
                                    tid_counter += 1
                                    texts.append({"textId": tid, "text": text})
                                    text_info.append({
                                        "textId": tid,
                                        "slideIndex": s_idx,
                                        "shapePath": path,
                                        "tableRowIndex": row_idx,
                                        "tableCellIndex": cell_idx,
                                        "paragraphIndex": p_idx,
                                        "runIndex": r_idx
                                    })

    _write_json(out_texts_json, texts)
    _write_json(out_meta_json, {
        "originalFilePath": os.path.abspath(input_path),
        "textInfo": text_info
    })
    return out_meta_json, out_texts_json

def importTextsToPptx(meta_json: str) -> None:
    texts_json = meta_json.replace("_meta.json", "_content.json")
    meta = _read_json(meta_json)
    texts_map = _texts_to_map(_read_json(texts_json))
    input_path = meta["originalFilePath"]

    prs = Presentation(input_path)
    for info in meta.get("textInfo", []):
        tid = info["textId"]
        if tid not in texts_map:
            continue
        s_idx = int(info["slideIndex"])
        shape_path = info["shapePath"]
        p_idx = int(info.get("paragraphIndex", -1))
        r_idx = int(info.get("runIndex", -1))
        table_row_idx = int(info.get("tableRowIndex", -1))
        table_cell_idx = int(info.get("tableCellIndex", -1))
        if s_idx < 0 or s_idx >= len(prs.slides):
            continue
        slide = prs.slides[s_idx]
        shape = _get_shape_by_path(slide, shape_path)
        if shape is None:
            continue
        # Check if the shape is a table
        if shape.has_table:
            if table_row_idx < 0 or table_row_idx >= len(shape.table.rows):
                continue
            if table_cell_idx < 0 or table_cell_idx >= len(shape.table.rows[table_row_idx].cells):
                continue
            cell = shape.table.cell(table_row_idx, table_cell_idx)
            if p_idx < 0 or p_idx >= len(cell.text_frame.paragraphs):
                continue
            paragraph = cell.text_frame.paragraphs[p_idx]
            if r_idx < 0 or r_idx >= len(paragraph.runs):
                continue
            run = paragraph.runs[r_idx]
            run.text = texts_map[tid]
        else:
            # If not a table, handle text frames
            if not (hasattr(shape, "has_text_frame") and shape.has_text_frame):
                continue
            tf = shape.text_frame
            if p_idx < 0 or p_idx >= len(tf.paragraphs):
                continue
            paragraph = tf.paragraphs[p_idx]
            if r_idx < 0 or r_idx >= len(paragraph.runs):
                continue
            run = paragraph.runs[r_idx]
            run.text = texts_map[tid]
    outputFileName = add_suffix_to_filename(input_path, "_Pimage")
    prs.save(outputFileName)
    return outputFileName


# testImagePath = os.path.join(currentDir, "source", "test.png")
# testImagePath_output = os.path.join(currentDir, "source", "test111.png")
testPdfPath = os.path.join(currentDir, "source", "xinqianli.pdf")
# testPdfPath_output = os.path.join(currentDir, "source", "图纸111.pdf")
testPptPath = os.path.join(currentDir, "source", "PDF图文.pptx")
# testPptPath_output = os.path.join(currentDir, "source", "PDF图文111.pptx")
testXslxPath = os.path.join(currentDir, "source", "sheet.xlsx")


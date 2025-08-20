from PIL import Image
import fitz
import os
import sys
import io
from pptx import Presentation
from pptx.util import Inches
import json
import numpy as np
import hashlib
from openpyxl import load_workbook
from openpyxl.drawing.image import Image as ExcelImage

currentDir = ""
if getattr(sys, 'frozen', False):
    currentDir = os.path.dirname(sys.executable)
else:
    currentDir = os.path.dirname(os.path.abspath(__file__))

def add_suffix_to_filename(file_path, suffix):
    # Split the file path into directory, base name, and extension
    directory, filename = os.path.split(file_path)
    name, ext = os.path.splitext(filename)
    # Create the new filename with the suffix
    new_filename = f"{name}{suffix}{ext}"
    # Join the new filename with the directory
    new_file_path = os.path.join(directory, new_filename)
    return new_file_path

def convert_numpy_types(data):
    if isinstance(data, dict):
        return {key: convert_numpy_types(value) for key, value in data.items()}
    elif isinstance(data, list):
        return [convert_numpy_types(item) for item in data]
    elif isinstance(data, np.integer):
        return int(data)  # Convert NumPy integer to Python int
    elif isinstance(data, np.floating):
        return float(data)  # Convert NumPy float to Python float
    else:
        return data  # Return the data as is if it's not a NumPy type

def extractImagesFromPDF(inputPath):
    # Load the PDF
    pdf_document = fitz.open(inputPath)
    inputFileName, inputFileExtension = os.path.splitext(os.path.basename(inputPath))
    extractionFolderPath = os.path.join(os.path.dirname(inputPath), f"{inputFileName}_extraction")
    os.makedirs(extractionFolderPath, exist_ok=True)
    extractionJsonPath = os.path.join(extractionFolderPath, f"{inputFileName}_images_meta.json")
    attributeSaveList = []
    image_hash_map = {}  # To store image hash and corresponding file paths
    # Iterate over all pages in the PDF
    for page_index in range(len(pdf_document)):
        page = pdf_document[page_index]
        images = page.get_images(full=True)  # Get all images on the page
        for image_index, img in enumerate(images):
            xref = img[0]  # XREF of the image
            base_image = pdf_document.extract_image(xref)  # Extract image data
            image_bytes = base_image["image"]
            image_stream = io.BytesIO(image_bytes)
            image = Image.open(image_stream)
            # Compute a hash of the image for comparison
            image_hash = hashlib.md5(image.tobytes()).hexdigest()
            # Check if the image is already saved
            if image_hash in image_hash_map:
                filePath = image_hash_map[image_hash]  # Use the existing file path
            else:
                # Save the image if it's not already saved
                filePath = os.path.join(extractionFolderPath, f"{inputFileName}_page{page_index}_img{image_index}.png")
                image.save(filePath)
                image_hash_map[image_hash] = filePath  # Store the hash and file path
            attributes = {
                "filePath": filePath,
                "pageIndex": page_index,
                "imageIndex": image_index,
                "imageWidth": base_image.get("width"),
                "imageHeight": base_image.get("height")
            }
            attributeSaveList.append(attributes)
    # Save attributes to a JSON file
    with open(extractionJsonPath, 'w', encoding='utf-8') as json_file:
        jsonData = {"originalFilePath": inputPath, "imageInfo":attributeSaveList}
        json.dump(jsonData, json_file, ensure_ascii=False, indent=2)
    return extractionJsonPath


def importImagesToPDF(jsonPath):
    # Load the JSON data
    with open(jsonPath, 'r', encoding='utf-8') as json_file:
        jsonData = json.load(json_file)
    originalFilePath = jsonData['originalFilePath']
    image_info = jsonData['imageInfo']
    # Open the original PDF document
    pdf_document = fitz.open(originalFilePath)
    # Iterate through the image info and add images back to the PDF
    for image_attributes in image_info:
        page_index = image_attributes['pageIndex']
        image_index = image_attributes['imageIndex']
        filePath = image_attributes['filePath']
        image_width = image_attributes['imageWidth']
        image_height = image_attributes['imageHeight']
        # Ensure the page exists in the PDF
        if page_index < len(pdf_document):
            page = pdf_document[page_index]
            # Get all images on the page
            images = page.get_images(full=True)
            # Check if the image index is valid
            if image_index < len(images):
                # Get the xref of the image to replace
                xref_to_replace = images[image_index][0]
                image_rects = page.get_image_rects(xref_to_replace)
                if image_rects:  # Check if the image was found on the page
                    original_rect = image_rects[0]  # Typically the first rectangle if image appears once
                    # print(f"Original image position: {original_rect}")
                    # First remove the existing image
                    page.delete_image(xref_to_replace)
                    # Insert the new image at the same position
                    page.insert_image(original_rect, filename=filePath)
                else:
                    print("Image not found on page")
    # Save the modified PDF to the specified output path
    outputFileName = add_suffix_to_filename(originalFilePath, "_Pimage")
    try:
        pdf_document.save(outputFileName, incremental=False)
    except ValueError:
        pdf_document.save(outputFileName, incremental=True)
    pdf_document.close()
    return outputFileName

# prtivate method for pptx
def extract_images_from_shapes(shapes, shapes_to_modify, page_index):
    for shape_index, shape in enumerate(shapes):
        # print(f"page_index:{page_index}/{shape_index}-{shape.shape_type}")
        if shape.shape_type == 13:  # Shape type 13 is for pictures
            shapes_to_modify.append((shape_index, shape))  # Store index and shape
        elif hasattr(shape, 'shapes'): 
            # Recursively extract images from grouped shapes
            extract_images_from_shapes(shape.shapes, shapes_to_modify, page_index)


def extractImagesFromPPTX(inputPath):
    # Load the presentation
    presentation = Presentation(inputPath)
    inputFileName, inputFileExtension = os.path.splitext(os.path.basename(inputPath))
    extractionFolderPath = os.path.join(os.path.dirname(inputPath), f"{inputFileName}_extraction")
    os.makedirs(extractionFolderPath, exist_ok=True)
    extractionJsonPath = os.path.join(extractionFolderPath, f"{inputFileName}_images_meta.json")
    attributeSaveList = []
    image_hash_map = {}  # To store image hash and corresponding file paths
    for page_index, slide in enumerate(presentation.slides):
        shapes_to_modify = []
        # Collect shapes to modify
        extract_images_from_shapes(slide.shapes, shapes_to_modify, page_index)
        # Modify collected shapes
        for shape_index, shape in shapes_to_modify:
            # Extract the image
            try:
                image_stream = io.BytesIO(shape.image.blob)
                image = Image.open(image_stream)
            except:
                print("skipped an image that is not loadable")
                continue
            # Compute a hash of the image for comparison
            image_hash = hashlib.md5(image.tobytes()).hexdigest()
            # Check if the image is already saved
            if image_hash in image_hash_map:
                filePath = image_hash_map[image_hash]  # Use the existing file path
            else:
                # Save the image if it's not already saved
                filePath = os.path.join(extractionFolderPath, f"{inputFileName}_{page_index}_{shape_index}.png")
                image.save(filePath)
                image_hash_map[image_hash] = filePath  # Store the hash and file path
            # Add shape attributes to the list
            attributes = {
                "filePath": filePath,
                "pageIndex": page_index,
                "shapeIndex": shape_index,
                "shapePosition": {"top": shape.top, "left": shape.left, "width": shape.width, "height": shape.height},
            }
            attributeSaveList.append(attributes)
    # Save attributes to a JSON file
    with open(extractionJsonPath, 'w', encoding='utf-8') as json_file:
        attributeSaveList = convert_numpy_types(attributeSaveList)
        jsonData = {"originalFilePath": inputPath, "imageInfo": attributeSaveList}
        json.dump(jsonData, json_file, ensure_ascii=False, indent=2)
    return extractionJsonPath

def importImagesToPPTX(sourceJson):
    with open(sourceJson, 'r', encoding='utf-8') as json_file:
        result = json.load(json_file)
    data = result["imageInfo"]
    pptxFile = result["originalFilePath"]
    presentation = Presentation(pptxFile)
    print(f"importing images to {pptxFile}")
    for entry in data:
        imagePath = entry["filePath"]
        if not os.path.isfile(imagePath):
            continue
        pageIndex = entry["pageIndex"]
        shapePosition = entry['shapePosition']
        slide = presentation.slides[pageIndex]
        slide.shapes.add_picture(imagePath, shapePosition["left"], shapePosition["top"], shapePosition["width"], shapePosition["height"])
    outputFileName = add_suffix_to_filename(pptxFile, "_Pimage")
    presentation.save(outputFileName)
    return outputFileName


def extractImagesFromXLSX(inputPath):
    # Load the workbook
    workbook = load_workbook(inputPath)
    inputFileName, inputFileExtension = os.path.splitext(os.path.basename(inputPath))
    extractionFolderPath = os.path.join(os.path.dirname(inputPath), f"{inputFileName}_extraction")
    os.makedirs(extractionFolderPath, exist_ok=True)
    extractionJsonPath = os.path.join(extractionFolderPath, f"{inputFileName}_images_meta.json")
    attributeSaveList = []
    image_hash_map = {}  # To store image hash and corresponding file paths

    for sheet_index, sheet in enumerate(workbook.sheetnames):
        ws = workbook[sheet]
        for image_index, img in enumerate(ws._images):
            # Extract the image
            if isinstance(img, ExcelImage):
                # Create a BytesIO object from the image
                image_stream = io.BytesIO(img._data())  # Use img._data() to get the image bytes
                image = Image.open(image_stream)
                # Compute a hash of the image for comparison
                image_hash = hashlib.md5(image.tobytes()).hexdigest()
                # Check if the image is already saved
                if image_hash in image_hash_map:
                    filePath = image_hash_map[image_hash]  # Use the existing file path
                else:
                    # Save the image if it's not already saved
                    filePath = os.path.join(extractionFolderPath, f"{inputFileName}_sheet{sheet_index}_img{image_index}.png")
                    image.save(filePath)
                    image_hash_map[image_hash] = filePath  # Store the hash and file path

                # Add image attributes to the list
                attributes = {
                    "filePath": filePath,
                    "sheetIndex": sheet_index,
                    "imageIndex": image_index,
                    "imagePosition": {"top": img.anchor._from.row, "left": img.anchor._from.col, "width": image.width, "height": image.height},
                }
                attributeSaveList.append(attributes)

    # Save attributes to a JSON file
    with open(extractionJsonPath, 'w', encoding='utf-8') as json_file:
        jsonData = {"originalFilePath": inputPath, "imageInfo": attributeSaveList}
        json.dump(jsonData, json_file, ensure_ascii=False, indent=2)
    
    return extractionJsonPath

def importImagesToXLSX(sourceJson):
    with open(sourceJson, 'r', encoding='utf-8') as json_file:
        result = json.load(json_file)
    data = result["imageInfo"]
    xlsxFile = result["originalFilePath"]
    workbook = load_workbook(xlsxFile)
    print(f"Importing images to {xlsxFile}")
    
    for entry in data:
        imagePath = entry["filePath"]
        if not os.path.isfile(imagePath):
            continue
        sheetIndex = entry["sheetIndex"]
        imagePosition = entry['imagePosition']
        ws = workbook.worksheets[sheetIndex]
        # Add the image to the worksheet
        img = ExcelImage(imagePath)
        if entry["imagePosition"]["width"] is not None:
            img.width = entry["imagePosition"]["width"]
            img.height = entry["imagePosition"]["height"]
        # Calculate the cell position
        row = imagePosition["top"] + 1  # Excel is 1-indexed
        col = imagePosition["left"] + 1  # Excel is 1-indexed
        # Convert column index to letter
        column_letter = ''
        while col > 0:
            col, remainder = divmod(col - 1, 26)
            column_letter = chr(remainder + ord('A')) + column_letter
        
        img.anchor = f"{column_letter}{row}"  # Set the anchor to the cell in A1 notation
        
        ws.add_image(img)
    outputFileName = add_suffix_to_filename(xlsxFile, "_Pimage")
    workbook.save(outputFileName)
    return outputFileName





# testImagePath = os.path.join(currentDir, "source", "test.png")
# testImagePath_output = os.path.join(currentDir, "source", "test111.png")
# testPdfPath = os.path.join(currentDir, "source", "图纸.pdf")
# testPdfPath_output = os.path.join(currentDir, "source", "图纸111.pdf")
# testPptPath = os.path.join(currentDir, "source", "PDF图文.pptx")
# testPptPath_output = os.path.join(currentDir, "source", "PDF图文111.pptx")
# testXslxPath = os.path.join(currentDir, "source", "sheet.xlsx")

# jsonPath = extractImageFromPPTX(testPptPath)
# importImageToPPTX(jsonPath)

# jsonPath = extractImagesFromPDF(testPdfPath)
# print(jsonPath)
# importImagesToPDF("c:\\Users\\Administrator\\Desktop\\tools\\others\\MoonTranslator\\source\\图纸_extraction\\图纸.pdf.json")

# extractImagesFromXLSX(testXslxPath)
# importImagesToXLSX("c:\\Users\\Administrator\\Desktop\\tools\\others\\MoonTranslator\\source\\sheet_extraction\\sheet.xlsx.json")